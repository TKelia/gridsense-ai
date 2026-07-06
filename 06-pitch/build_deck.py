"""Builds the GridSense AI capstone defense deck (.pptx) from the verified script.
Dark theme matching the app. On-slide text is concise; full speaker notes go into
each slide's notes pane. Run: python build_deck.py"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(HERE)
OUT = os.path.join(HERE, "GridSense-AI-Defense-Deck.pptx")

# palette
BG     = RGBColor(0x0B, 0x0F, 0x17)
PANEL  = RGBColor(0x16, 0x1C, 0x2B)
EMER   = RGBColor(0x34, 0xD3, 0x99)
CYAN   = RGBColor(0x22, 0xD3, 0xEE)
AMBER  = RGBColor(0xFB, 0xBF, 0x24)
ROSE   = RGBColor(0xF4, 0x3F, 0x5E)
WHITE  = RGBColor(0xF1, 0xF5, 0xF9)
MUTED  = RGBColor(0x94, 0xA3, 0xB8)
DIM    = RGBColor(0x64, 0x74, 0x8B)
DARKTX = RGBColor(0x0B, 0x0F, 0x17)
FONT   = "Segoe UI"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    f = s.background.fill
    f.solid(); f.fore_color.rgb = BG
    return s


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


def box(s, x, y, w, h):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def setrun(r, text, size, color=WHITE, bold=False, italic=False):
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = FONT


def para(tf, first=False):
    return tf.paragraphs[0] if first else tf.add_paragraph()


def rich(p, segments, size, base=WHITE):
    """segments: list of (text, bold, color)"""
    for txt, bold, col in segments:
        r = p.add_run()
        setrun(r, txt, size, col or base, bold)


def parse_bold(text):
    """split on *bold* markers -> list of (text, bold)"""
    out, parts = [], text.split("*")
    for i, part in enumerate(parts):
        if part:
            out.append((part, i % 2 == 1))
    return out


def title(s, text, kicker=None):
    # accent rule
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.55), Inches(0.12), Inches(0.62))
    bar.fill.solid(); bar.fill.fore_color.rgb = EMER; bar.line.fill.background()
    tf = box(s, 0.85, 0.45, 11.8, 1.0)
    p = para(tf, True)
    setrun(p.add_run(), text, 30, WHITE, True)
    if kicker:
        p2 = para(tf)
        setrun(p2.add_run(), kicker, 13, EMER, False, True)


def footer(s, source, num):
    tf = box(s, 0.6, 6.95, 9.0, 0.4)
    setrun(para(tf, True).add_run(), source, 9, DIM, italic=True)
    tn = box(s, 9.4, 6.95, 3.7, 0.4)
    p = para(tn, True); p.alignment = PP_ALIGN.RIGHT
    setrun(p.add_run(), "gridsense-ai-zeta.vercel.app", 9, DIM)


def bullets(s, items, x=0.95, y=1.65, w=11.6, h=4.9, size=17, gap=10):
    tf = box(s, x, y, w, h)
    for i, it in enumerate(items):
        p = para(tf, i == 0)
        p.space_after = Pt(gap)
        # bullet dot
        d = p.add_run(); setrun(d, "●  ", size, EMER, True)
        for txt, bold in parse_bold(it):
            r = p.add_run(); setrun(r, txt, size, WHITE, bold)
    return tf


def tier_bar(s, x, y, w):
    """mini tier gauge like the app: 89 / 310 / 369"""
    segs = [("89", EMER, 0.30), ("310", AMBER, 0.34), ("369", ROSE, 0.36)]
    cx = x
    for label, col, frac in segs:
        sw = w * frac
        rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(cx), Inches(y), Inches(sw), Inches(0.5))
        rect.fill.solid(); rect.fill.fore_color.rgb = col; rect.line.fill.background()
        tfr = rect.text_frame; tfr.word_wrap = True
        pr = tfr.paragraphs[0]; pr.alignment = PP_ALIGN.CENTER
        setrun(pr.add_run(), label + " RWF", 12, DARKTX, True)
        cx += sw
    cap = box(s, x, y + 0.55, w, 0.35)
    setrun(para(cap, True).add_run(), "0–20 kWh        20–50 kWh              above 50 kWh", 10, MUTED)


def table(s, rows, x, y, w, col_w, header=True, fs=12, rh=0.0):
    nr, nc = len(rows), len(rows[0])
    gt = s.shapes.add_table(nr, nc, Inches(x), Inches(y), Inches(w), Inches(0.5 * nr)).table
    gt.first_row = header
    gt.horz_banding = False
    for ci, cw in enumerate(col_w):
        gt.columns[ci].width = Inches(cw)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            c = gt.cell(ri, ci)
            c.margin_left = Inches(0.12); c.margin_right = Inches(0.1)
            c.margin_top = Inches(0.05); c.margin_bottom = Inches(0.05)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.fill.solid()
            c.fill.fore_color.rgb = PANEL if (ri == 0 and header) else BG
            tf = c.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            ish = (ri == 0 and header)
            for txt, bold in parse_bold(str(val)):
                r = p.add_run()
                setrun(r, txt, fs, EMER if ish else WHITE, bold or ish)
    return gt


def arch_node(s, x, y, w, h, head, sub):
    n = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    n.fill.solid(); n.fill.fore_color.rgb = PANEL
    n.line.color.rgb = EMER; n.line.width = Pt(1.5)
    n.shadow.inherit = False
    tf = n.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    setrun(p.add_run(), head, 14, WHITE, True)
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    setrun(p2.add_run(), sub, 10.5, MUTED)


def arch_arrow(s, x, y, w, h=0.42):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    a.fill.solid(); a.fill.fore_color.rgb = EMER; a.line.fill.background()
    a.shadow.inherit = False


def draw_architecture(s, y=3.95):
    nodes = [
        ("Sensors", "CT clamp + smart plugs"),
        ("ESP32", "reads current · Wi-Fi"),
        ("Cloud", "/api/ingest (JSON)"),
        ("Dashboard", "+ tariff-aware AI"),
    ]
    bw, bh, aw, gap = 2.4, 1.35, 0.5, 0.16
    total = len(nodes) * bw + (len(nodes) - 1) * (aw + 2 * gap)
    x = (13.333 - total) / 2
    for i, (head, sub) in enumerate(nodes):
        arch_node(s, x, y, bw, bh, head, sub)
        x += bw
        if i < len(nodes) - 1:
            arch_arrow(s, x + gap, y + (bh - 0.42) / 2, aw)
            x += aw + 2 * gap
    cap = box(s, 0.95, y + bh + 0.35, 11.4, 0.7)
    pc = para(cap, True); pc.alignment = PP_ALIGN.CENTER
    rich(pc, [("The simulated demo emits the ", False, MUTED), ("same JSON contract", True, EMER),
              (" — real hardware swaps in with a one-line change.", False, MUTED)], 13)


def payback_chart(s, x, y, w, h):
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x - 0.12), Inches(y - 0.12),
                              Inches(w + 0.24), Inches(h + 0.24))
    card.fill.solid(); card.fill.fore_color.rgb = WHITE; card.line.fill.background()
    card.shadow.inherit = False
    data = CategoryChartData()
    data.categories = ["10% savings", "15% savings", "20% savings"]
    data.add_series("Core monitor", (13.6, 9.1, 6.8))
    data.add_series("Full kit", (23.0, 15.6, 11.7))
    gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(x), Inches(y), Inches(w), Inches(h), data)
    ch = gf.chart
    ch.has_title = True
    ch.chart_title.text_frame.text = "Months to payback — 100 kWh home (lower is better)"
    ch.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
    ch.chart_title.text_frame.paragraphs[0].runs[0].font.bold = True
    ch.chart_title.text_frame.paragraphs[0].runs[0].font.color.rgb = DARKTX
    ch.font.size = Pt(11); ch.font.color.rgb = DARKTX; ch.font.name = FONT
    ch.has_legend = True
    ch.legend.position = XL_LEGEND_POSITION.BOTTOM
    ch.legend.include_in_layout = False
    plot = ch.plots[0]
    plot.gap_width = 80
    plot.has_data_labels = True
    plot.data_labels.number_format = '0.0'
    plot.data_labels.number_format_is_linked = False
    plot.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    plot.data_labels.font.size = Pt(9); plot.data_labels.font.color.rgb = DARKTX
    colors = [EMER, RGBColor(0x38, 0xBD, 0xF8)]
    for i, series in enumerate(plot.series):
        series.format.fill.solid(); series.format.fill.fore_color.rgb = colors[i]
        series.format.line.fill.background()
    va = ch.value_axis
    va.has_major_gridlines = True
    va.major_gridlines.format.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    va.tick_labels.font.size = Pt(9); va.tick_labels.font.color.rgb = DARKTX
    ca = ch.category_axis
    ca.tick_labels.font.size = Pt(10); ca.tick_labels.font.color.rgb = DARKTX


def img_or_placeholder(s, path, x, y, h):
    if os.path.exists(path):
        return s.shapes.add_picture(path, Inches(x), Inches(y), height=Inches(h))
    ph = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(h * 1.1), Inches(h))
    ph.fill.solid(); ph.fill.fore_color.rgb = PANEL; ph.line.color.rgb = DIM
    setrun(ph.text_frame.paragraphs[0].add_run(), "[screenshot]", 12, MUTED)
    return ph


# ============================================================ SLIDE 1 — Title
s = slide()
_logo = os.path.join(ROOT, "05-build", "dashboard", "public", "icon-512.png")
if os.path.exists(_logo):
    s.shapes.add_picture(_logo, Inches(5.8665), Inches(1.45), height=Inches(1.6))
else:
    logo = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.85), Inches(1.45), Inches(1.6), Inches(1.6))
    logo.fill.solid(); logo.fill.fore_color.rgb = EMER; logo.line.fill.background()
    lp = logo.text_frame.paragraphs[0]; lp.alignment = PP_ALIGN.CENTER
    setrun(lp.add_run(), "G", 60, DARKTX, True)
tf = box(s, 1.0, 3.35, 11.3, 2.7)
p = para(tf, True); p.alignment = PP_ALIGN.CENTER
setrun(p.add_run(), "GridSense AI", 48, WHITE, True)
p2 = para(tf); p2.alignment = PP_ALIGN.CENTER
setrun(p2.add_run(), "See your power. Spend less.", 18, EMER, italic=True)
pp = para(tf); pp.alignment = PP_ALIGN.CENTER
setrun(pp.add_run(), "An energy-intelligence platform for Rwanda — homes · landlords · enterprises", 14, WHITE)
purl = para(tf); purl.alignment = PP_ALIGN.CENTER
setrun(purl.add_run(), "https://gridsense-ai-zeta.vercel.app", 14, CYAN, True)
p3 = para(tf); p3.alignment = PP_ALIGN.CENTER; p3.space_before = Pt(14)
setrun(p3.add_run(), "Tesi Songa Kelia  ·  BSc Software Engineering  ·  African Leadership University, Kigali", 13, MUTED)
p4 = para(tf); p4.alignment = PP_ALIGN.CENTER
setrun(p4.add_run(), "Capstone Defense  ·  June 2026", 12, DIM)
notes(s, "Open with confidence. 'Rwandans just got a more expensive, tiered electricity bill — and almost no way to see it coming. GridSense AI fixes that.' The tagline is the whole thesis.")

# ============================================================ SLIDE 2 — Problem
s = slide()
title(s, "The problem", "You can't manage what you can't see")
bullets(s, [
    "Oct 2025: residential power became steeply *tiered* — *89 → 310 → 369 RWF/kWh*. Average tariff up *~15%*.",
    "Households pay via prepaid *\"cash power\"* — they see a falling balance, not which appliance drains it.",
    "No warning when a home crosses from a cheap tier into an expensive one.",
], y=1.6, h=2.4)
tier_bar(s, 0.95, 4.45, 8.0)
tf = box(s, 0.95, 5.55, 11.4, 1.1)
rich(para(tf, True), [("Creep from 20 → 21 kWh and the marginal unit costs ", False, WHITE),
                      ("3.5× more", True, AMBER), (" — with zero visibility into why.", False, WHITE)], 16)
footer(s, "Source: 01-research §1 — REG tariffs, RURA press release", "2")
notes(s, "Make them feel it. The tiered cliff is the pain. A home creeping one kWh over a threshold pays 3.5x more on the marginal unit and has no idea. This invisibility is what we sell against.")

# ============================================================ SLIDE 3 — Insight
s = slide()
title(s, "The insight that reshaped our strategy", "We confront the brutal truth instead of hiding it")
bullets(s, [
    "*~50%* of connected homes use *≤ 20 kWh/month* (the 89 RWF lifeline). A 20% saving there ≈ *356 RWF (~$0.25)/mo*. No one buys a device for that.",
    "The money — and willingness to pay — is with homes, landlords & SMEs that *cross into the 310 & 369 RWF tiers*.",
    "*So we don't target \"every household.\" We target the tier-crossers.*",
], y=1.7, h=3.6, gap=14)
footer(s, "Source: 01-research §3; 02-strategy §A, §D1", "3")
notes(s, "This is the rigor slide that wins points. We confronted the weakness (half of homes can't save meaningfully) and it reshaped targeting. Pre-empts the panel's #1 question: who actually pays, and does the math work?")

# ============================================================ SLIDE 4 — Solution
s = slide()
title(s, "The solution")
bullets(s, [
    "A low-cost kit: *whole-home CT clamp + ESP32 (Wi-Fi)* for real-time total use & live RWF cost, *plus energy-metering smart plugs* for the biggest appliances.",
    "A simple *web app + tariff-aware AI*: live cost, bill forecast, *tier-cliff alerts*, and specific savings advice.",
    "Localized: *RWF, Rwanda's real tiers, Kinyarwanda-ready*.",
], y=1.7, h=3.8, gap=14)
footer(s, "Source: DIRECTION.md; 02-strategy §D2", "4")
notes(s, "Cheap, safe, self-installable. CT = whole-home picture; plugs = 'which appliance?'. The intelligence is the localization and the tariff-aware AI, not exotic hardware.")

# ============================================================ SLIDE 4b — One platform, three customers
s = slide()
title(s, "One platform — three customers", "Same verified tariff engine powers all three")
bullets(s, [
    "*Home (free):* a household sees its own live cost, tier-cliff alerts & personalized savings.",
    "*Landlord / Property:* manage many units & tenants — *exact per-tenant bills*, invoices, reminders, get paid on time.",
    "*Enterprise:* hotels & large compounds — many properties, teams, exports, connected devices.",
], y=1.75, h=3.6, gap=16)
footer(s, "Live: choose a workspace after sign-in · 05-build (workspace.tsx)", "4b")
notes(s, "GridSense isn't just a home gadget — it's a platform. After sign-in you pick a workspace. The same exact tariff engine serves a single home, a landlord with many tenants, and an enterprise. Landlords are the clearest paying wedge.")

# ============================================================ SLIDE 4c — Landlord billing
s = slide()
title(s, "Landlords: bill fairly, get paid on time", "The feature that ends \"I don't know how I'm charged\"")
bullets(s, [
    "Add *properties → units → tenants* (name, phone, email) + meter readings.",
    "*Exact amount due* per tenant from the real tariff — with billing period and due date.",
    "One tap: *Invoice / Receipt / Consumption Report* (printable) → *send via WhatsApp / SMS / email*.",
    "*Payment reminders* as a calendar (.ics): 5 / 3 / 2 / 1 days before, the morning of, and 12:00 on the due day.",
    "Transparent, itemized charges that respect tenant privacy (Law 058/2021).",
], y=1.7, h=4.4, gap=12)
footer(s, "Live: Properties workspace · 05-build (landlord.tsx, Documents.tsx, billing.ts)", "4c")
notes(s, "This is the monetizable headline. Tenants constantly dispute power charges; GridSense shows the exact math, generates the invoice, sends it, and schedules the reminders. One buyer (the landlord), many meters, recurring value.")

# ============================================================ SLIDE 4d — Pricing
s = slide()
title(s, "Plans & pricing", "Individuals free for adoption; landlords & enterprise pay")
table(s, [
    ["Plan", "Price (indicative)", "For"],
    ["Home", "Free", "A single household"],
    ["Landlord — Starter", "~RWF 9,900/mo  (≤5 units)", "Small landlords"],
    ["Landlord — Growth", "~RWF 29,000/mo  (≤20 units)", "Growing portfolios"],
    ["Enterprise", "from ~RWF 150,000/mo", "Hotels, estates, teams"],
    ["Device add-on", "kit 40,100–69,100 + ~2,500/device/mo", "Automatic live data"],
], x=0.95, y=1.8, w=11.4, col_w=[3.2, 4.6, 3.6], fs=13)
tf = box(s, 0.95, 5.5, 11.4, 1.0)
rich(para(tf, True), [("Premium checkout + cart are live; ", False, WHITE),
                      ("demo mode charges no card", True, EMER),
                      (" and is payment-provider-ready (MTN MoMo / Flutterwave / Stripe).", False, WHITE)], 13)
footer(s, "Indicative pricing — 02-strategy/platform-expansion.md §E", "4d")
notes(s, "Pricing is indicative and labelled as such. Home is free to drive adoption and data; revenue comes from landlords and enterprise plus the device add-on. Checkout is a real premium flow but charges nothing in demo — integration-ready, never faked.")

# ============================================================ SLIDE 5 — Demo
s = slide()
title(s, "Live demo", "English & Kinyarwanda — same app, one toggle")
left = img_or_placeholder(s, os.path.join(ROOT, "gridsense-livenow.png"), 0.7, 1.5, 5.0)
right = img_or_placeholder(s, os.path.join(ROOT, "gridsense-rw-live.png"), 7.0, 1.5, 5.0)
gap = Inches(0.4)
total = left.width + gap + right.width
startx = int((SW - total) / 2)
left.left = startx; left.top = Inches(1.5)
right.left = startx + left.width + gap; right.top = Inches(1.5)
cap = box(s, 0.6, 6.55, 12.1, 0.5)
pc = para(cap, True); pc.alignment = PP_ALIGN.CENTER
rich(pc, [("4 screens — Live Now · This Month · Appliances · Save.   ", False, WHITE),
          ("Sensor data is simulated & labelled; the tariff math is real.", True, EMER)], 13)
footer(s, "Built & tested: 9 unit + 5 e2e tests · Lighthouse 100/100/100/100 · 05-build/", "5")
notes(s, "DO THE LIVE DEMO HERE — open the Vercel URL. Walk Live Now, trigger the tier-cliff alert, show This Month forecast, Appliances, Save. SWITCH TO KINYARWANDA LIVE — it lands. Say out loud: 'sensor data is simulated and labelled; the tariff math and appliance model are real and sourced.' Honesty is a strength.")

# ============================================================ SLIDE 6 — AI honest
s = slide()
title(s, "Is the \"AI\" real? — told honestly")
bullets(s, [
    "*Now:* tariff-aware analytics (kWh → live RWF across real tiers) + *rule-based* recommendations + run-rate *forecasting*.",
    "*Next:* ML bill forecasting + anomaly detection (\"your fridge is drawing 30% more than usual\").",
    "*Later:* appliance disaggregation (*NILM*) as data grows.",
    "*We never claim NILM we haven't built.*",
], y=1.7, h=3.9, gap=13)
footer(s, "Source: 02-strategy §D3 — visible in the Save screen + roadmap", "6")
notes(s, "Pre-empts 'is the AI real?'. Honest scoping is a defense, not a weakness. A panel respects a clear roadmap over an over-claim it can puncture.")

# ============================================================ SLIDE 7 — Architecture
s = slide()
title(s, "How it works — and why it's buildable")
bullets(s, [
    "Non-invasive *CT clamp on the main feed*; *ESP32* reads current and uploads over Wi-Fi.",
    "Built on *proven open-source* designs (ESP32 + SCT-013 + EmonLib; CircuitSetup).",
], y=1.6, h=1.7, gap=10)
draw_architecture(s, y=3.95)
footer(s, "Source: 01-research §5; 05-build/dashboard-build-plan.md, src/lib/types.ts", "7")
notes(s, "Pre-empts 'did you really build hardware?'. Hardware is off-the-shelf and proven; our innovation is localization + tariff-aware AI + GTM. The real ingestion contract makes the demo a genuine product slice, not a mockup.")

# ============================================================ SLIDE 8 — BOM
s = slide()
title(s, "Feasibility & cost — sourced, in RWF")
table(s, [
    ["Component", "Local price (RWF)", "Source"],
    ["ESP32 dev board", "15,500", "SoftTech Supply (Kigali)"],
    ["SCT-013-100 CT clamp", "12,000", "SoftTech Supply"],
    ["Metering smart plug", "14,500", "SoftTech Supply"],
    ["PSU + wiring", "7,600", "SoftTech Supply"],
    ["Whole-home monitor (entry SKU)", "~40,100  (~$27)", "core unit"],
    ["Full kit (+2 plugs)", "~69,100  (~$47)", "1 USD ≈ 1,464 RWF"],
], x=0.95, y=1.75, w=11.4, col_w=[5.0, 3.4, 3.0], fs=14)
footer(s, "Source: 04-business/capex-opex.md §1 — priced from live Kigali retailers (taxes embedded)", "8")
notes(s, "We priced locally in RWF, not off a foreign catalogue — the most defensible, conservative basis (bulk import would be cheaper). 93% of the kit is live-sourced; the small remainder is flagged.")

# ============================================================ SLIDE 9 — Payback
s = slide()
title(s, "Unit economics & payback — the math is real")
tf = box(s, 0.95, 1.5, 11.5, 0.6)
rich(para(tf, True), [("Verified tiers → a ", False, WHITE), ("100 kWh home pays 29,530 RWF/month", True, EMER),
                      (".  Payback = kit cost ÷ monthly RWF saved.", False, WHITE)], 15)
payback_chart(s, 0.85, 2.35, 7.0, 4.2)
bullets(s, [
    "Core monitor: *~9 months @ 15%* savings.",
    "Heavy users (200 kWh): *~6 months @ 10%*.",
    "*< 12 months* for tier-crossers ✓",
    "Shown as a *range* — we don't fake one number.",
], x=8.35, y=2.55, w=4.5, h=4.0, size=15, gap=14)
footer(s, "Source: 04-business/capex-opex.md §4 — tariff engine unit-tested in 05-build", "9")
notes(s, "Bills computed straight from the verified RURA tiers (tested in code: 100→29,530, 150→47,980, 200→66,430). Payback under 12 months for tier-crossers proves the beachhead economically. The savings % is a labelled assumption — see Risks.")

# ============================================================ SLIDE 10 — Market
s = slide()
title(s, "Market & beachhead")
bullets(s, [
    "*1.946M* households on the grid; *84.6%* national electricity access.",
    "Beachhead = *tier-crossing urban homes + landlords/compounds + SMEs* — they cross 310/369, own appliances, have Wi-Fi + smartphones, and feel the hike.",
    "Mass-market & lifeline homes = a later *social-impact tier* (grant-funded), not the paying entry.",
], y=1.7, h=3.8, gap=14)
footer(s, "Source: 01-research §2–3; 02-strategy §D1", "10")
notes(s, "Landlords are arguably the best wedge — one buyer, many meters, a billing-dispute pain point. SMEs have a clear ROI line. Mass-market is impact, not the paying entry.")

# ============================================================ SLIDE 11 — Competition
s = slide()
title(s, "Competition & our moat")
bullets(s, [
    "Sense (~$300, US, ML) · Emporia Vue · IoTaWatt (DIY) · Bidgely (utility B2B).",
    "*None* is localized: no RWF-tier logic, no Kinyarwanda, not Rwanda-priced or -serviced.",
    "*Moat = localization + tariff-aware AI + go-to-market*, at a Rwandan price.",
], y=1.7, h=3.6, gap=14)
footer(s, "Source: 01-research §4", "11")
notes(s, "Strong tech exists — we don't pretend to out-engineer Sense. We win on the wedge they ignore: a tariff-aware, Kinyarwanda, affordably-serviced product for Rwanda.")

# ============================================================ SLIDE 12 — Regulation
s = slide()
title(s, "Regulation & compliance — a credibility asset")
bullets(s, [
    "*Behind-the-meter, non-invasive* — no utility permit to clip a CT on your own wiring; *we never touch the REG/EUCL meter*.",
    "Board-side install by a *licensed electrician* (REG/RURA requirement); smart plugs are plug-and-play.",
    "*Data:* compliant with *Law 058/2021*; register with *NCSA*; consent + encryption + data minimization.",
], y=1.7, h=3.8, gap=14)
footer(s, "Source: 01-research §6 — RURA, REG, RwandaLII (verified, not assumed)", "12")
notes(s, "Pre-empts 'will it mess with the meter?' and 'what about privacy?'. We verified the rule rather than assuming, and built compliance in from day one.")

# ============================================================ SLIDE 13 — GTM/Funding
s = slide()
title(s, "Go-to-market & funding — the lean ladder")
bullets(s, [
    "*$0 to win the capstone* (free cloud, open-source, ALU resources). Funding is for scaling.",
    "Ladder: demo → *Hanga PitchFest* (50M RWF; RDB+MINICT+UNDP) → 3–5 kits + 1 pilot → larger grants (FEC 2026, Youth4Climate $30k, ClimaFii $70k, EEP Africa €200–500k) → production.",
    "Channel: direct + *landlord/SME* partnerships + possible REG partnership.",
], y=1.7, h=3.9, gap=13)
footer(s, "Source: 04-business/fundraising-plan.md", "13")
notes(s, "Each rung needs only the proof from the one below — never a big upfront ask. Hanga is the #1 target (2026 call not yet published — we're prototype-ready and watching).")

# ============================================================ SLIDE 14 — Roadmap
s = slide()
title(s, "Roadmap")
bullets(s, [
    "*Phase 1 (now):* whole-home CT + ESP32 (Wi-Fi) + plugs; tariff-aware analytics, alerts, rules; web app. (demoed today)",
    "*Phase 2:* multi-circuit \"by-area\" CT in the board (partner electrician); GSM option; landlord sub-metering.",
    "*Phase 3:* ML forecasting → anomaly detection → NILM disaggregation.",
], y=1.7, h=3.8, gap=14)
footer(s, "Source: DIRECTION.md; 02-strategy §D2", "14")
notes(s, "Clear, staged, honest. Each phase unlocks the next as data and funding grow.")

# ============================================================ SLIDE 15 — Risks
s = slide()
title(s, "Risks & honest mitigations")
table(s, [
    ["Risk", "Mitigation"],
    ["Savings % is an assumption (we show a 5–20% range)", "Pilot will measure it; cite behavioral-feedback literature"],
    ["Hardware funding (~$27–47 / kit)", "$0 capstone path; bench PoC if free parts; Hanga/grants to scale"],
    ["Wi-Fi penetration in target homes", "Beachhead chosen for Wi-Fi/smartphone ownership; GSM in Phase 2"],
    ["Hanga 2026 dates unconfirmed", "Monitoring official channels; prototype-ready now"],
], x=0.95, y=1.8, w=11.4, col_w=[5.4, 6.0], fs=13)
footer(s, "Source: open questions in PROJECT-REPORT.md, 01-research, 04-business", "15")
notes(s, "Showing the gaps WITH mitigations beats pretending there are none. This is the slide that earns trust. Be calm and direct here.")

# ============================================================ SLIDE 16 — Close
s = slide()
title(s, "Impact & the ask")
bullets(s, [
    "*Impact:* energy literacy, lower bills, less waste, reduced peak demand, climate — real and fundable.",
    "*The ask:* capstone validation · a first pilot home/landlord · a path into Hanga PitchFest.",
], y=1.7, h=2.6, gap=16)
tf = box(s, 1.0, 4.7, 11.3, 1.2)
p = para(tf, True); p.alignment = PP_ALIGN.CENTER
setrun(p.add_run(), "GridSense AI — see your power, spend less.", 26, EMER, True)
footer(s, "Thank you — questions welcome (Q&A defense rehearsed)", "16")
notes(s, "Close on the mission and a concrete ask — Tesi's strength. Land it with conviction. Thank the panel; invite questions. You're ready (appendix).")

# ============================================================ SLIDE 17 — Q&A appendix
s = slide()
title(s, "Appendix — Q&A defense", "Rehearse these cold")
table(s, [
    ["If they ask…", "Answer"],
    ["Who actually pays?", "Tier-crossers: urban homes, landlords, SMEs. Economics on slides 3, 9, 10."],
    ["Sense/Emporia exist.", "None localized / RWF-tier-aware / Kinyarwanda / Rwanda-priced (11)."],
    ["Is the AI real?", "Tariff-aware + rules now; ML next; NILM later — no over-claim (6)."],
    ["Did you build hardware?", "Proven off-the-shelf; real ingestion contract; live demo (5, 7)."],
    ["Will it touch the meter?", "No — behind-the-meter, non-invasive; licensed electrician for board (12)."],
    ["Privacy?", "Law 058/2021; NCSA registration; consent + encryption (12)."],
    ["How do you make money?", "Device sale + optional subscription; landlord/SME tiers (13)."],
    ["Where did the numbers come from?", "Live Kigali retailers (RWF) + verified RURA tiers; savings as a range (8, 9)."],
], x=0.7, y=1.7, w=12.0, col_w=[3.8, 8.2], fs=12)
footer(s, "Full attack list: 02-strategy §E · sources in 01-research, 04-business", "17")
notes(s, "Don't show this unless asked — it's your cheat sheet. Each answer maps to a slide you can jump to.")

prs.save(OUT)
print("Saved:", OUT)
print("Slides:", len(prs.slides._sldIdLst))
